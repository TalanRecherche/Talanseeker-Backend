"""Created on Wed Aug 30 14:27:23 2023

@author: agarc

"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.cv_information_retrieval.abc_reader import ABCReader


class PPTXReader(ABCReader):
    @staticmethod
    def read_text(file_path: str) -> str | None:
        """Read all text from a PowerPoint (PPTX) file.

        Args:
        ----
        file_path (str): The path of the PPTX file.

        Returns:
        -------
        str: The extracted text from the PPTX file.

        Dependencies:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        """
        # load presentation
        prs = Presentation(file_path)
        # instantiate empty output
        text_runs = []
        # loop through slides and slide.shapes
        for slide in prs.slides:
            for shape in slide.shapes:
                # get text from shape
                shape_text = PPTXReader._extract_text_from_shape(shape)
                if shape_text:
                    text_runs.append(shape_text)

        # remove special characters
        text_runs = [
            run.encode("latin-1", errors="ignore").decode("latin-1")
            for run in text_runs
        ]

        # remove trailing spaces
        text_runs = [run.strip() for run in text_runs]

        # generate full text
        text = "\n\n\n".join(text_runs)
        # triple spaces are mapped to line breaks
        text = text.replace("   ", "\n")
        # remove multiple spaces
        text = text.replace("  ", " ")

        # remove trailing spaces and linebreaks
        text = text.rstrip("\n")
        text = text.strip()

        # return only if not empty
        if text:
            return text
        else:
            return None

    @staticmethod
    def _extract_text_from_shape(shape: Presentation) -> str:
        """Extract any text from multiple shapes types
        This could be defined recursively also but this seems easier that way.
        """
        sub_texts = []
        if shape.has_text_frame:  # from text frame
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    sub_texts.append(run.text)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # from tables
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            sub_texts.append(run.text)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # from groups
            for sub_shape in shape.shapes:
                sub_texts.extend(PPTXReader._extract_text_from_shape(sub_shape))
        elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    sub_texts.append(run.text)

        return " ".join(sub_texts).strip()
